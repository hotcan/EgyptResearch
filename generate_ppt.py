#!/usr/bin/env python3
"""Generate a PowerPoint presentation for the Egypt travel journal."""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))

# Colors
GOLD = RGBColor(0xC8, 0x97, 0x3A)
GOLD_DIM = RGBColor(0x8A, 0x6A, 0x28)
INK = RGBColor(0xE8, 0xE0, 0xD0)
INK_MUTED = RGBColor(0x8A, 0x7F, 0x6E)
BG = RGBColor(0x14, 0x12, 0x0E)
BG_SURFACE = RGBColor(0x1C, 0x18, 0x13)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Fonts
CN_FONT = "PingFang SC"
EN_FONT = "Georgia"

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=INK, bold=False, font_name=CN_FONT, alignment=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=18, color=INK, bold=False,
             font_name=CN_FONT, alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_gold_line(slide, left, top, width):
    """Add a thin gold horizontal line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


def add_image_fit(slide, img_path, left, top, max_w, max_h):
    """Add image scaled to fit within max_w x max_h."""
    from PIL import Image
    try:
        with Image.open(img_path) as im:
            iw, ih = im.size
    except Exception:
        return None

    ratio = min(max_w / iw, max_h / ih)
    w = int(iw * ratio)
    h = int(ih * ratio)
    # Center within the area
    x = left + (max_w - w) // 2
    y = top + (max_h - h) // 2
    return slide.shapes.add_picture(img_path, x, y, w, h)


def make_cover(prs):
    """Slide 1: Cover page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Eyebrow
    add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
                "TRAVEL JOURNAL  ·  2026", font_size=14, color=GOLD,
                font_name=EN_FONT)

    # Title
    add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.5),
                "埃及旅行手记", font_size=54, color=INK, bold=True)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(3.8), Inches(8), Inches(1),
                '2026年春节，随\u201c哲学课堂\u201d一起，沿着开罗、亚历山大、卢克索和沙姆沙伊赫，\n探寻世界宗教的起源，记录下思考和收获。',
                font_size=18, color=INK_MUTED)

    # Route
    add_gold_line(slide, Inches(1.5), Inches(5.3), Inches(10))
    add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
                "中国 ✈ 开罗 → 亚历山大 → 开罗 ✈ 卢克索 ✈ 沙姆沙伊赫 ✈ 开罗 ✈ 中国",
                font_size=16, color=GOLD, font_name=CN_FONT)

    # Date range
    add_textbox(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
                "2026.02.22 — 2026.03.04  ·  11天  ·  4城市  ·  5000年文明跨度",
                font_size=13, color=INK_MUTED, font_name=EN_FONT)


def make_route_overview(prs):
    """Slide 2: Route overview with stats."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.5),
                "ROUTE OVERVIEW", font_size=14, color=GOLD, font_name=EN_FONT)

    add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
                "行程路线", font_size=36, color=INK, bold=True)

    # Route with city descriptions
    cities = [
        ("开罗 Cairo", "Day 1-5, 9-11", "萨拉丁城堡 · 吉萨金字塔 · 大埃及博物馆 · 老馆 · 穆伊兹街"),
        ("亚历山大 Alexandria", "Day 4", "亚历山大图书馆 · 卡特贝城堡 · 地中海"),
        ("卢克索 Luxor", "Day 5-7", "帝王谷 · 卡纳克神庙 · 卢克索神庙 · 尼罗河"),
        ("沙姆沙伊赫 Sharm", "Day 7-9", "红海浮潜 · 四季酒店 · 西奈半岛"),
    ]

    y = Inches(2.8)
    for city, days, desc in cities:
        add_textbox(slide, Inches(1.5), y, Inches(4), Inches(0.4),
                    city, font_size=20, color=GOLD, bold=True)
        add_textbox(slide, Inches(5.5), y, Inches(2), Inches(0.4),
                    days, font_size=14, color=INK_MUTED, font_name=EN_FONT)
        add_textbox(slide, Inches(1.5), y + Inches(0.4), Inches(10), Inches(0.4),
                    desc, font_size=14, color=INK_MUTED)
        y += Inches(1.1)


def make_day_title_slide(prs, day_num, date, title, route, desc):
    """Create a day title slide with text info."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Day number eyebrow
    eyebrow = f"DAY {day_num:02d}  ·  {date}"
    add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.5),
                eyebrow, font_size=14, color=GOLD, font_name=EN_FONT)

    # Title (strip HTML tags)
    clean_title = title.replace("<em>", "").replace("</em>", "")
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
                clean_title, font_size=44, color=INK, bold=True)

    # Route
    add_textbox(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(0.5),
                route, font_size=16, color=GOLD)

    # Gold line
    add_gold_line(slide, Inches(1.5), Inches(4.2), Inches(4))

    # Description
    add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(2.5),
                desc, font_size=16, color=INK_MUTED)


def make_photo_slide(prs, photos, day_num, caption=""):
    """Create a slide with 1-3 photos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SURFACE)

    n = len(photos)
    valid_photos = [p for p in photos if os.path.exists(p)]
    if not valid_photos:
        add_textbox(slide, Inches(3), Inches(3), Inches(7), Inches(1),
                    f"Day {day_num} Photos", font_size=24, color=INK_MUTED,
                    alignment=PP_ALIGN.CENTER)
        return

    n = len(valid_photos)

    if n == 1:
        add_image_fit(slide, valid_photos[0],
                      Inches(1), Inches(0.5),
                      Inches(11.333), Inches(6))
    elif n == 2:
        add_image_fit(slide, valid_photos[0],
                      Inches(0.3), Inches(0.5),
                      Inches(6.2), Inches(6))
        add_image_fit(slide, valid_photos[1],
                      Inches(6.8), Inches(0.5),
                      Inches(6.2), Inches(6))
    elif n >= 3:
        # Left: large photo
        add_image_fit(slide, valid_photos[0],
                      Inches(0.3), Inches(0.3),
                      Inches(7), Inches(6.5))
        # Right top
        add_image_fit(slide, valid_photos[1],
                      Inches(7.6), Inches(0.3),
                      Inches(5.4), Inches(3.1))
        # Right bottom
        add_image_fit(slide, valid_photos[2],
                      Inches(7.6), Inches(3.7),
                      Inches(5.4), Inches(3.1))

    if caption:
        add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                    caption, font_size=11, color=INK_MUTED,
                    alignment=PP_ALIGN.CENTER)


def make_summary_slide(prs, section_num, title, bullets):
    """Create a summary section slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Section number
    add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.5),
                f"SUMMARY  ·  {section_num}", font_size=14, color=GOLD,
                font_name=EN_FONT)

    # Title
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
                title, font_size=36, color=INK, bold=True)

    add_gold_line(slide, Inches(1.5), Inches(3.2), Inches(3))

    # Bullets
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = INK_MUTED
        p.font.name = CN_FONT
        p.space_before = Pt(10)
        p.line_spacing = Pt(28)


def make_ending(prs):
    """Final slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
                "埃及没有给我答案\n——它给了我更好的问题",
                font_size=36, color=INK, bold=True, alignment=PP_ALIGN.CENTER)

    add_gold_line(slide, Inches(4), Inches(4.5), Inches(5))

    add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
                "EGYPT  ·  MMXXVI", font_size=16, color=GOLD,
                font_name=EN_FONT, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
                "hotcan  ·  2026年春节", font_size=14, color=INK_MUTED,
                alignment=PP_ALIGN.CENTER)


def main():
    # Load days.json
    with open(os.path.join(BASE, "data", "days.json"), "r", encoding="utf-8") as f:
        raw = json.load(f)
    days_data = raw.get("days", raw) if isinstance(raw, dict) else raw

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- COVER ---
    make_cover(prs)

    # --- ROUTE OVERVIEW ---
    make_route_overview(prs)

    # --- DAY SLIDES ---
    # Photo selections per day
    photo_map = {
        1: [
            os.path.join(BASE, "assets/day1/muhammad_ali_mosque.jpg"),
            os.path.join(BASE, "assets/day1/cairo_overview.jpg"),
            os.path.join(BASE, "assets/day1/neferdidi_tea.jpg"),
        ],
        2: [
            os.path.join(BASE, "assets/day2/giza_pyramids.jpg"),
            os.path.join(BASE, "assets/day2/sphinx.jpg"),
            os.path.join(BASE, "assets/day2/papyrus_tutankhamun.jpg"),
        ],
        3: [
            os.path.join(BASE, "days/day3/action/IMG_9791.jpg"),
            os.path.join(BASE, "days/day3/action/IMG_9890.jpg"),
            os.path.join(BASE, "days/day3/action/IMG_9931.jpg"),
        ],
        4: [
            os.path.join(BASE, "days/day4/action/reading_room.jpg"),
            os.path.join(BASE, "days/day4/action/entrance.jpg"),
            os.path.join(BASE, "days/day4/action/citadel.jpg"),
        ],
        5: [
            os.path.join(BASE, "days/day5/history/gem_exterior.jpg"),
            os.path.join(BASE, "days/day5/history/tutankhamun_throne.jpg"),
            os.path.join(BASE, "days/day5/history/solar_boat_wide.jpg"),
        ],
        6: [
            os.path.join(BASE, "days/day6/photos/seti_ceiling.jpg"),
            os.path.join(BASE, "days/day6/photos/twelve_baboons.jpg"),
            os.path.join(BASE, "days/day6/photos/felucca_sunset.jpg"),
        ],
        7: [
            os.path.join(BASE, "days/day7/photos/IMG_1088.jpg"),
            os.path.join(BASE, "days/day7/photos/IMG_1305.jpg"),
            os.path.join(BASE, "days/day7/photos/IMG_1407.jpg"),
        ],
        8: [
            os.path.join(BASE, "days/day8/photos/IMG_1524.jpg"),
            os.path.join(BASE, "days/day8/photos/IMG_1682.jpg"),
            os.path.join(BASE, "days/day8/photos/IMG_1772.jpg"),
        ],
        9: [],  # No photos
        10: [
            os.path.join(BASE, "days/day10/photos/IMG_1840.jpg"),
            os.path.join(BASE, "days/day10/photos/IMG_2036.jpg"),
            os.path.join(BASE, "days/day10/photos/IMG_2241.jpg"),
        ],
        11: [],  # No photos
    }

    photo_captions = {
        1: "穆罕默德·阿里清真寺 · 城堡俯瞰开罗 · 娜芙缇缇奶茶店",
        2: "吉萨金字塔群 · 狮身人面像 · 图坦卡蒙莎草纸画",
        3: "悬空教堂 · 圣安东尼修道院 · 攀登千级台阶",
        4: "亚历山大图书馆主阅览室 · 图书馆入口 · 卡特贝城堡",
        5: "大埃及博物馆 · 图坦卡蒙黄金宝座 · 胡夫太阳船",
        6: "塞提一世墓天花板 · 图坦卡蒙墓12只狒狒 · 尼罗河帆船日落",
        7: "卡纳克巨柱大厅 · 卢克索神庙 · 夜景灯光",
        8: "红海浮潜 · 四季酒店 · 海滨度假",
        10: "埃及博物馆老馆 · 祖韦伊拉门 · 穆伊兹街",
    }

    # Day descriptions for title slides (short versions)
    day_highlights = {
        1: "第一次踏上非洲大陆。在萨拉丁城堡俯瞰整个开罗，参观穆罕默德·阿里清真寺，在城堡里一家温州人开的奶茶店喝了第一杯「全球化」。",
        2: "上午在吉萨高地面对三座大金字塔和狮身人面像，下午在哈利利市场带回四幅手绘莎草纸画。晚宴上与开罗大学教授畅谈哲学。",
        3: "上午探访科普特开罗的悬空教堂，下午穿越三小时沙漠抵达圣安东尼修道院。攀登一千多级台阶到达修道士山洞，思考苦修的意义。",
        4: "驱车三小时到亚历山大图书馆，获得 VIP 参观。午餐鱼市场海鲜。返程大巴上从斐洛的寓意解经聊到加缪的荒谬。",
        5: "参观大埃及博物馆（GEM）——图坦卡蒙宝藏、胡夫太阳船、拉美西斯巨像。傍晚飞往卢克索，开启上埃及之旅。",
        6: "上午卢克索博物馆，下午帝王谷：塞提一世最美陵墓、图坦卡蒙墓的12只狒狒壁画。傍晚尼罗河帆船看日落。",
        7: "卡纳克134根巨柱撑起的大厅、卢克索神庙的日落与夜灯。阿布·哈加格清真寺直接建在法老柱廊之上——三千年被垂直压缩。",
        8: "从黄沙切换到碧蓝。红海浮潜、四季酒店沙滩。距以色列边境200公里的度假胜地，手机弹出美以联合空袭伊朗的新闻。",
        9: "下午从沙姆沙伊赫飞回开罗。飞机上读了大半本古埃及历史。晚上在如意坊吃了一顿中餐。轻松的一天。",
        10: "独自探索开罗。埃及博物馆老馆、伊斯兰艺术博物馆、祖韦伊拉门、穆伊兹街漫步，最后在费沙维咖啡馆喝了一杯薄荷茶。",
        11: "十一天的旅程画上句号。从金字塔到穆伊兹街，从亡灵书到薄荷茶——带着塞满照片和想法的脑袋，从开罗飞回上海。",
    }

    for i, day in enumerate(days_data):
        day_num = i + 1
        date = day.get("date", "")
        title = day.get("title", day.get("card_title", f"Day {day_num}"))
        route = day.get("route", "")
        desc = day_highlights.get(day_num, day.get("desc", ""))

        # Title slide for this day
        make_day_title_slide(prs, day_num, date, title, route, desc)

        # Photo slide (if photos exist)
        photos = photo_map.get(day_num, [])
        if photos:
            caption = photo_captions.get(day_num, "")
            make_photo_slide(prs, photos, day_num, caption)

    # --- SUMMARY SLIDES ---
    make_summary_slide(prs, "I", "三层历史", [
        "古埃及三千年：从纳尔迈统一到托勒密——金字塔、神庙、亡灵书的实物印证",
        "亚伯拉罕三宗教两千年：犹太教、基督教、伊斯兰教同源分流",
        "伊斯兰教在埃及一千四百年：从阿拉伯征服到穆伊兹街的马穆鲁克建筑",
        "三层历史叠加在同一片土地上——这就是埃及",
    ])

    make_summary_slide(prs, "II", "建筑的奇迹", [
        "金字塔：230万块石头垒起146米——四千五百年前没有铁器、车轮、起重机",
        "卡纳克134根巨柱：法老们压倒一切的野心",
        "帝王谷：三千年前的壁画颜色依然鲜艳得像昨天画的",
        "穆伊兹街：从石头到木头，从法老到苏丹，建筑语言变了，对不朽的追求没变",
    ])

    make_summary_slide(prs, "III", "蓝色与红色", [
        "前七天是密集的博物馆、神庙、陵墓、沙漠——信息过载",
        "第八天飞到沙姆沙伊赫，从黄沙切换到碧蓝——什么都不用想的一天",
        "但蓝色背后藏着红色：距以色列边境200公里，手机弹出美以空袭伊朗",
        "和平与战争可以在同一个时空里平行运行——这是沙姆沙伊赫教我的",
    ])

    make_summary_slide(prs, "IV", "哲学的沉淀", [
        "圣安东尼的千级台阶：肉体的极限是大脑停止噪音的方式",
        "亡灵书：死亡不是终结，而是另一段旅程的开始",
        "卡纳克巨柱 vs 紫禁城金柱：石头与木头，永恒与更替",
        "「明知会消逝仍要创造」——这种顽固，就是人之为人最动人的部分",
    ])

    # --- ENDING ---
    make_ending(prs)

    # Save
    output = os.path.join(BASE, "egypt-trip.pptx")
    prs.save(output)
    print(f"PPT saved to: {output}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
